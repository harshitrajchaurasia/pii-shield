"""
PI Remover File Processor Module.

Contains file processing functions for various formats:
- CSV (with multiprocessing support)
- Excel (.xlsx, .xls)
- JSON
- Text (.txt, .md, .log)
- Word (.docx)
- PowerPoint (.pptx)
- PDF
- HTML/XML

Usage:
    from pi_remover.processors.file_processor import (
        process_file,
        process_csv,
        process_dataframe,
    )
"""

import json
import logging
import multiprocessing as mp
from datetime import datetime
from pathlib import Path
from typing import Any, List, Optional, Tuple

import pandas as pd
from tqdm import tqdm

from pi_remover.config import PIRemoverConfig, config_to_dict, config_from_dict
from pi_remover.utils import (
    IS_WINDOWS,
    get_multiprocessing_method,
)

# Initialize logger
logger = logging.getLogger("pi_remover")


# Optional Imports
# Security validation (optional - for file processing)
try:
    from pi_remover.security import validate_file_security
    SECURITY_AVAILABLE = True
except ImportError:
    SECURITY_AVAILABLE = False
    def validate_file_security(path: str, strict: bool = True) -> Tuple[bool, Optional[str]]:  # type: ignore[misc]
        return True, None  # No-op if security module not available

# Resource monitoring for auto-scaling (optional)
try:
    import sys
    import os
    sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))))
    from shared.resource_monitor import ResourceMonitor
    RESOURCE_MONITOR_AVAILABLE = True
except ImportError:
    RESOURCE_MONITOR_AVAILABLE = False
    ResourceMonitor = None  # type: ignore


# Validation
class ValidationError(Exception):
    """Input validation failed."""
    pass


def validate_file(file_path: str, must_exist: bool = True) -> Path:
    """
    Validate a file path.
    
    Args:
        file_path: Path to validate
        must_exist: Whether the file must exist
        
    Returns:
        Path object
        
    Raises:
        ValidationError: If validation fails
    """
    if not file_path:
        raise ValidationError("File path cannot be empty")
    
    path = Path(file_path)
    
    if must_exist and not path.exists():
        raise ValidationError(f"File not found: {file_path}")
    
    if must_exist and not path.is_file():
        raise ValidationError(f"Not a file: {file_path}")
    
    # Check supported extensions
    supported_extensions = {'.csv', '.xlsx', '.xls', '.json', '.txt', '.md', '.log',
                           '.docx', '.doc', '.pptx', '.ppt', '.pdf', '.html', '.htm', '.xml'}
    if path.suffix.lower() not in supported_extensions:
        raise ValidationError(
            f"Unsupported file type: {path.suffix}. "
            f"Supported: {', '.join(sorted(supported_extensions))}"
        )
    
    return path


def validate_columns(df: pd.DataFrame, columns: List[str]) -> List[str]:
    """
    Validate that columns exist in DataFrame.
    
    Args:
        df: DataFrame to check
        columns: List of column names
        
    Returns:
        List of valid columns (with warnings for missing ones)
    """
    valid_columns = []
    for col in columns:
        if col in df.columns:
            valid_columns.append(col)
        else:
            logger.warning(f"Column '{col}' not found in DataFrame. Skipping.")
    
    if not valid_columns:
        raise ValidationError(
            f"None of the specified columns found. "
            f"Available columns: {list(df.columns)}"
        )
    
    return valid_columns


def validate_config(config: PIRemoverConfig) -> PIRemoverConfig:
    """
    Validate configuration settings.
    
    Args:
        config: Configuration to validate
        
    Returns:
        Validated config (possibly with corrections)
    """
    # Ensure at least one detection method is enabled
    if not any([config.enable_ner, config.enable_regex, config.enable_dictionaries]):
        logger.warning("All detection methods disabled. Enabling regex.")
        config.enable_regex = True
    
    # Ensure at least one PI type is enabled
    pi_types = [
        config.redact_names, config.redact_emails, config.redact_phones,
        config.redact_emp_ids, config.redact_asset_ids, config.redact_ip_addresses,
        config.redact_urls, config.redact_hostnames, config.redact_credentials
    ]
    if not any(pi_types):
        logger.warning("All PI types disabled. Enabling common types.")
        config.redact_emails = True
        config.redact_phones = True
        config.redact_names = True
    
    # Validate worker count
    if config.num_workers < 1:
        config.num_workers = 1
    elif config.num_workers > mp.cpu_count() * 2:
        logger.warning(f"num_workers ({config.num_workers}) exceeds recommended max. Setting to {mp.cpu_count()}.")
        config.num_workers = mp.cpu_count()
    
    # Validate batch size
    if config.batch_size < 100:
        logger.warning("batch_size too small. Setting to 100.")
        config.batch_size = 100
    
    return config


# Multiprocessing Workers
# Global worker remover (initialized in worker processes)
_worker_remover = None


def _init_worker(config_dict: dict):
    """Initialize worker process with PIRemover instance."""
    global _worker_remover
    # Import here to avoid circular imports
    from pi_remover.remover import PIRemover
    
    # Reconstruct config from dict
    config = config_from_dict(config_dict)
    config.show_progress = False  # Disable progress in workers
    config.use_multiprocessing = False  # Already in worker
    
    _worker_remover = PIRemover(config)


def _process_chunk_worker(chunk_data: Tuple[int, pd.DataFrame, List[str]]) -> pd.DataFrame:
    """Worker function to process a chunk of data."""
    global _worker_remover
    idx, chunk, columns = chunk_data

    # Use safe per-row application to avoid single-row failures
    assert _worker_remover is not None, "Worker remover not initialized"
    for col in columns:
        if col in chunk.columns:
            try:
                chunk[f"{col}_cleaned"] = _safe_apply_series(
                    chunk[col], _worker_remover, _worker_remover.config, col
                )
            except Exception:
                # If worker-level error and cannot continue, re-raise
                raise

    return chunk


def _safe_apply_series(series: pd.Series, remover, config: PIRemoverConfig, col_name: str) -> pd.Series:
    """
    Apply redaction to a pandas Series with per-row try/except and logging.

    Args:
        series: Input pandas Series
        remover: PIRemover instance
        config: Configuration
        col_name: Column name for logging
        
    Returns:
        New Series with same index
    """
    results = []
    error_count = 0

    iterator = series.items()
    if getattr(config, 'show_progress', False):
        iterator = tqdm(series.items(), desc=f"Redacting {col_name}")

    for idx, val in iterator:
        try:
            if pd.isna(val):
                results.append(val)
            else:
                results.append(remover.redact(str(val)))
        except Exception as e:
            error_count += 1
            # Log if requested
            try:
                if config.error_log_file:
                    with open(config.error_log_file, 'a', encoding='utf-8') as f:
                        f.write(f"{datetime.now().isoformat()} | Column:{col_name} | Index:{idx} | Error:{str(e)}\n")
                        if config.include_original_in_log:
                            f.write(f"Original: {repr(val)}\n")
            except Exception:
                # Swallow logging errors to avoid cascading failures
                pass

            if config.continue_on_error:
                results.append(val)
            else:
                raise

        if config.max_errors and error_count >= config.max_errors > 0:
            raise RuntimeError("Maximum error threshold reached.")

    return pd.Series(results, index=series.index)


# Data Processing Functions
def process_dataframe(
    df: pd.DataFrame, 
    columns: List[str],
    config: Optional[PIRemoverConfig] = None,
    suffix: str = "_cleaned"
) -> pd.DataFrame:
    """
    Process a DataFrame, redacting PI from specified columns.

    Args:
        df: Input DataFrame
        columns: List of column names to process
        config: PIRemoverConfig instance
        suffix: Suffix for cleaned column names

    Returns:
        DataFrame with new cleaned columns added
    """
    # Import here to avoid circular imports
    from pi_remover.remover import PIRemover
    
    config = config or PIRemoverConfig()
    config = validate_config(config)
    columns = validate_columns(df, columns)
    
    remover = PIRemover(config)
    result = df.copy()

    for col in columns:
        new_col = f"{col}{suffix}"
        logger.info(f"Processing column: {col} -> {new_col}")

        # Use safe per-row application to prevent single-row failures
        result[new_col] = _safe_apply_series(result[col], remover, config, col)

    return result


def process_csv(
    input_path: str, 
    output_path: str, 
    columns: List[str],
    config: Optional[PIRemoverConfig] = None,
    chunksize: int = 10000
) -> None:
    """
    Process a large CSV file in chunks with optional multiprocessing.

    Args:
        input_path: Path to input CSV
        output_path: Path to output CSV
        columns: Columns to process
        config: PIRemoverConfig instance
        chunksize: Number of rows per chunk
    """
    config = config or PIRemoverConfig()

    # Determine total rows for progress bar with error handling
    try:
        total_rows = sum(1 for _ in open(input_path, 'r', encoding='utf-8', errors='ignore')) - 1
    except UnicodeDecodeError:
        logger.warning(f"UTF-8 decode failed for {input_path}, trying latin-1")
        try:
            total_rows = sum(1 for _ in open(input_path, 'r', encoding='latin-1', errors='ignore')) - 1
        except Exception as e:
            logger.error(f"Failed to read CSV {input_path}: {type(e).__name__}: {e}")
            raise
    except Exception as e:
        logger.error(f"Failed to read CSV {input_path}: {type(e).__name__}: {e}")
        raise

    # Auto-scaling: Determine optimal workers based on file size and resources
    if config.auto_scale_workers and RESOURCE_MONITOR_AVAILABLE:
        monitor = ResourceMonitor()
        profile = monitor.get_scaling_profile(total_rows)
        
        # Update config based on scaling profile
        config.num_workers = profile.num_workers
        chunksize = profile.chunk_size
        
        # Determine if multiprocessing should be used
        should_use_mp = (
            config.use_multiprocessing and 
            profile.num_workers > 1 and 
            total_rows >= config.multiprocessing_threshold
        )
        
        logger.info(f"Auto-scaling: {total_rows:,} rows -> {profile.name} profile "
                   f"({profile.num_workers} workers, chunk_size={chunksize})")
    else:
        # Manual mode: use config settings directly
        should_use_mp = (
            config.use_multiprocessing and 
            config.num_workers > 1 and
            total_rows >= config.multiprocessing_threshold
        )

    if should_use_mp:
        # Multiprocessing mode
        _process_csv_multiprocessing(input_path, output_path, columns, config, chunksize, total_rows)
    else:
        # Single-threaded mode
        _process_csv_single(input_path, output_path, columns, config, chunksize, total_rows)


def process_file(
    input_path: str, 
    output_path: str, 
    columns: List[str],
    config: Optional[PIRemoverConfig] = None, 
    chunksize: int = 10000,
    skip_security_check: bool = False
) -> None:
    """
    Process a file based on its extension.

    Supports: CSV, Excel, JSON, TXT, DOCX, PPTX, PDF, HTML, and more.
    
    Args:
        input_path: Path to input file
        output_path: Path to output file
        columns: Columns to process (for structured files)
        config: PIRemoverConfig instance
        chunksize: Chunk size for CSV processing
        skip_security_check: Skip file security validation (not recommended)
    """
    # Security check first
    if not skip_security_check:
        is_valid, error = validate_file_security(input_path)
        if not is_valid:
            logger.error(f"Security check failed for '{input_path}': {error}")
            raise ValueError(f"File rejected: {error}")
    
    ext = Path(input_path).suffix.lower()
    cfg = config or PIRemoverConfig()

    if ext in {'.xlsx', '.xls'}:
        _process_excel(input_path, output_path, columns, cfg)
    elif ext == '.json':
        _process_json(input_path, output_path, columns, cfg)
    elif ext in {'.txt', '.md', '.log'}:
        _process_txt(input_path, output_path, cfg)
    elif ext in {'.docx', '.doc'}:
        _process_docx(input_path, output_path, cfg)
    elif ext in {'.pptx', '.ppt'}:
        _process_pptx(input_path, output_path, cfg)
    elif ext == '.pdf':
        _process_pdf(input_path, output_path, cfg)
    elif ext in {'.html', '.htm', '.xml'}:
        _process_html(input_path, output_path, cfg)
    elif ext == '.csv':
        process_csv(input_path, output_path, columns, config=cfg, chunksize=chunksize)
    else:
        # Try as plain text for unknown extensions
        _process_txt(input_path, output_path, cfg)


# Internal File Processors
def _process_csv_single(
    input_path: str, 
    output_path: str, 
    columns: List[str],
    config: PIRemoverConfig, 
    chunksize: int, 
    total_rows: int
) -> None:
    """Single-threaded CSV processing."""
    from pi_remover.remover import PIRemover
    
    remover = PIRemover(config)

    first_chunk = True
    processed = 0

    with tqdm(total=total_rows, desc="Processing CSV (single-threaded)") as pbar:
        for chunk in pd.read_csv(input_path, chunksize=chunksize,
                                 encoding='utf-8', on_bad_lines='skip'):
            # Process each specified column
            for col in columns:
                if col in chunk.columns:
                    new_col = f"{col}_cleaned"
                    chunk[new_col] = _safe_apply_series(chunk[col], remover, config, col)

            # Write to output
            if first_chunk:
                chunk.to_csv(output_path, index=False, mode='w')
                first_chunk = False
            else:
                chunk.to_csv(output_path, index=False, mode='a', header=False)

            processed += len(chunk)
            pbar.update(len(chunk))

    logger.info(f"Processed {processed} rows -> {output_path}")


def _process_csv_multiprocessing(
    input_path: str, 
    output_path: str, 
    columns: List[str],
    config: PIRemoverConfig, 
    chunksize: int, 
    total_rows: int
) -> None:
    """
    Multiprocessing CSV processing with spaCy NER support.
    
    Platform considerations:
    - Windows uses 'spawn': Workers don't inherit globals, must use initializer
    - Linux uses 'fork': Workers copy parent process, but we still use initializer for consistency
    - maxtasksperchild: Recycles workers to prevent memory leaks (important for spaCy)
    """
    num_workers = config.num_workers
    config_dict = config_to_dict(config)
    
    # Get current start method for logging
    start_method = mp.get_start_method() or get_multiprocessing_method()
    
    logger.info(f"Loading data for multiprocessing ({num_workers} workers, {start_method} method)...")

    # Read entire file into memory for parallel processing
    df = pd.read_csv(input_path, encoding='utf-8', on_bad_lines='skip')
    logger.info(f"Loaded {len(df)} rows")

    # Split into chunks for parallel processing
    chunk_size = max(1000, len(df) // num_workers)
    chunks = [df[i:i+chunk_size].copy() for i in range(0, len(df), chunk_size)]
    chunk_data = [(i, chunk, columns) for i, chunk in enumerate(chunks)]

    logger.info(f"Processing with {num_workers} workers, {len(chunks)} chunks...")

    # Platform-aware Pool configuration:
    # - maxtasksperchild: Recycle workers to free memory (spaCy models can leak)
    maxtasks = 50 if IS_WINDOWS else 100  # More frequent recycling on Windows
    
    try:
        # Process in parallel with worker initialization
        with mp.Pool(
            processes=num_workers, 
            initializer=_init_worker, 
            initargs=(config_dict,),
            maxtasksperchild=maxtasks  # Prevents memory leaks from spaCy
        ) as pool:
            results = list(tqdm(
                pool.imap(_process_chunk_worker, chunk_data),
                total=len(chunks),
                desc=f"Processing CSV ({num_workers} workers)"
            ))

        # Combine results
        result_df = pd.concat(results, ignore_index=True)

        # Save
        result_df.to_csv(output_path, index=False)
        logger.info(f"Processed {len(result_df)} rows -> {output_path}")
        
    except Exception as e:
        # Fallback to single-threaded processing if multiprocessing fails
        logger.warning(f"Multiprocessing failed ({e}), falling back to single-threaded mode")
        _process_csv_single(input_path, output_path, columns, config, chunksize, total_rows)


def _process_excel(
    input_path: str, 
    output_path: str, 
    columns: List[str],
    config: PIRemoverConfig
) -> None:
    """Process Excel file."""
    try:
        df = pd.read_excel(input_path, sheet_name=0)
        logger.info(f"Loaded Excel file with {len(df)} rows")
    except Exception as e:
        logger.error(f"Failed to read Excel file '{input_path}': {e}")
        return

    result_df = process_dataframe(df, columns, config=config)

    try:
        result_df.to_excel(output_path, index=False)
        logger.info(f"Processed {len(result_df)} rows -> {output_path}")
    except PermissionError:
        logger.error(f"Permission denied writing to '{output_path}'")
    except OSError as e:
        logger.error(f"OS error writing to '{output_path}': {type(e).__name__}: {e}")
    except Exception as e:
        logger.error(f"Failed to write Excel output '{output_path}': {e}")


def _process_json(
    input_path: str, 
    output_path: str, 
    columns: List[str],
    config: PIRemoverConfig
) -> None:
    """
    Process JSON file (array of objects or single object).
    
    Supports:
    - Array of objects: [{"field1": "...", "field2": "..."}, ...]
    - Single object: {"field1": "...", "field2": "..."}
    - Nested objects: will process string values at any depth
    """
    from pi_remover.remover import PIRemover
    
    # Read JSON with fallback encodings
    try:
        with open(input_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        logger.info(f"Loaded JSON file: {input_path}")
    except json.JSONDecodeError as e:
        logger.error(f"Invalid JSON in '{input_path}': {e}")
        return
    except UnicodeDecodeError:
        logger.warning(f"UTF-8 decode failed for {input_path}, trying latin-1")
        try:
            with open(input_path, 'r', encoding='latin-1') as f:
                data = json.load(f)
            logger.info(f"Loaded JSON file with latin-1 encoding: {input_path}")
        except Exception as e:
            logger.error(f"Failed to read JSON file '{input_path}': {e}")
            return
    except Exception as e:
        logger.error(f"Failed to read JSON file '{input_path}': {e}")
        return

    remover = PIRemover(config)

    def redact_value(val: Any, field_path: str = "") -> Any:
        """Recursively redact string values in JSON structure."""
        try:
            if isinstance(val, str):
                return remover.redact(val)
            elif isinstance(val, dict):
                return {k: redact_value(v, f"{field_path}.{k}") for k, v in val.items()}
            elif isinstance(val, list):
                return [redact_value(item, f"{field_path}[{i}]") for i, item in enumerate(val)]
            else:
                return val  # Numbers, booleans, null stay unchanged
        except Exception as e:
            logger.warning(f"Failed to redact value at {field_path}: {type(e).__name__}")
            return val

    if isinstance(data, list):
        # Array of objects
        result = []
        for i, item in enumerate(tqdm(data, desc="Processing JSON")):
            if isinstance(item, dict):
                if columns:
                    # Only process specified columns
                    processed = item.copy()
                    for col in columns:
                        if col in processed and isinstance(processed[col], str):
                            processed[f"{col}_cleaned"] = remover.redact(processed[col])
                    result.append(processed)
                else:
                    # Process all string values
                    result.append(redact_value(item))
            else:
                result.append(redact_value(item))
        logger.info(f"Processed {len(result)} JSON records")
    else:
        # Single object
        if columns:
            result = data.copy()
            for col in columns:
                if col in result and isinstance(result[col], str):
                    result[f"{col}_cleaned"] = remover.redact(result[col])
        else:
            result = redact_value(data)
        logger.info("Processed single JSON object")

    try:
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(result, f, indent=2, ensure_ascii=False)
        logger.info(f"Saved JSON output -> {output_path}")
    except PermissionError:
        logger.error(f"Permission denied writing to '{output_path}'")
    except OSError as e:
        logger.error(f"OS error writing to '{output_path}': {type(e).__name__}: {e}")


def _process_txt(input_path: str, output_path: str, config: PIRemoverConfig) -> None:
    """
    Process plain text file (line by line) with per-line error handling.
    
    Each line is treated as independent text to redact.
    Failures on individual lines preserve the original line.
    """
    from pi_remover.remover import PIRemover
    
    try:
        with open(input_path, 'r', encoding='utf-8', errors='replace') as f:
            lines = f.readlines()
        logger.info(f"Loaded TXT file with {len(lines)} lines")
    except Exception as e:
        logger.error(f"Failed to read TXT file '{input_path}': {e}")
        return

    remover = PIRemover(config)
    
    result_lines = []
    for line in tqdm(lines, desc="Processing TXT"):
        # Preserve line endings, with per-line error handling
        stripped = line.rstrip('\n\r')
        try:
            redacted = remover.redact(stripped)
            result_lines.append(redacted + '\n')
        except Exception as e:
            logger.warning(f"Failed to redact line, keeping original: {type(e).__name__}")
            result_lines.append(line)

    try:
        with open(output_path, 'w', encoding='utf-8') as f:
            f.writelines(result_lines)
        logger.info(f"Processed {len(lines)} lines -> {output_path}")
    except PermissionError:
        logger.error(f"Permission denied writing to '{output_path}'")
    except OSError as e:
        logger.error(f"OS error writing to '{output_path}': {type(e).__name__}: {e}")


def _process_docx(input_path: str, output_path: str, config: PIRemoverConfig) -> None:
    """Process Word DOCX file - extracts text from paragraphs and tables."""
    try:
        from docx import Document
    except ImportError:
        logger.error("python-docx not installed. Run: pip install python-docx")
        return
    
    from pi_remover.remover import PIRemover
    
    try:
        doc = Document(input_path)
        logger.info(f"Loaded DOCX file: {input_path}")
    except Exception as e:
        logger.error(f"Failed to read DOCX file '{input_path}': {e}")
        return
    
    remover = PIRemover(config)
    
    # Process paragraphs
    for para in doc.paragraphs:
        if para.text.strip():
            para.text = remover.redact(para.text)
    
    # Process tables
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    cell.text = remover.redact(cell.text)
    
    try:
        doc.save(output_path)
        logger.info(f"Processed DOCX -> {output_path}")
    except Exception as e:
        logger.error(f"Failed to write DOCX output '{output_path}': {e}")


def _process_pptx(input_path: str, output_path: str, config: PIRemoverConfig) -> None:
    """Process PowerPoint PPTX file - extracts text from slides."""
    try:
        from pptx import Presentation
    except ImportError:
        logger.error("python-pptx not installed. Run: pip install python-pptx")
        return
    
    from pi_remover.remover import PIRemover
    
    try:
        prs = Presentation(input_path)
        logger.info(f"Loaded PPTX file with {len(prs.slides)} slides")
    except Exception as e:
        logger.error(f"Failed to read PPTX file '{input_path}': {e}")
        return
    
    remover = PIRemover(config)
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.text.strip():
                                run.text = remover.redact(run.text)
    
    try:
        prs.save(output_path)
        logger.info(f"Processed PPTX -> {output_path}")
    except Exception as e:
        logger.error(f"Failed to write PPTX output '{output_path}': {e}")


def _process_pdf(input_path: str, output_path: str, config: PIRemoverConfig) -> None:
    """
    Process PDF file - extracts text, redacts, saves as TXT.
    
    Note: PDFs can't be easily modified in-place, so output is plain text.
    """
    try:
        import pdfplumber
    except ImportError:
        logger.error("pdfplumber not installed. Run: pip install pdfplumber")
        return
    
    from pi_remover.remover import PIRemover
    
    # Additional PDF security check - scan for JavaScript/actions
    try:
        with open(input_path, 'rb') as f:
            pdf_bytes = f.read(50000)  # First 50KB
            pdf_str = pdf_bytes.decode('latin-1', errors='ignore')
            
            dangerous_patterns = ['/JavaScript', '/JS ', '/Launch', '/OpenAction', 
                                  '/AA ', '/URI ', '/GoToR', '/EmbeddedFile']
            for pattern in dangerous_patterns:
                if pattern in pdf_str:
                    logger.warning(f"PDF contains potentially dangerous content: {pattern}")
                    if SECURITY_AVAILABLE:
                        from pi_remover.security import SecurityConfig
                        if SecurityConfig.QUARANTINE_SUSPICIOUS_FILES:
                            logger.error(f"Rejecting PDF with {pattern}")
                            return
    except Exception as e:
        logger.warning(f"Could not scan PDF for dangerous content: {e}")
    
    try:
        with pdfplumber.open(input_path) as pdf:
            logger.info(f"Loaded PDF with {len(pdf.pages)} pages")
            
            remover = PIRemover(config)
            result_lines = []
            
            for i, page in enumerate(pdf.pages):
                text = page.extract_text() or ""
                if text.strip():
                    redacted = remover.redact(text)
                    result_lines.append(f"--- Page {i+1} ---\n{redacted}\n\n")
            
            # Save as text
            txt_output = output_path if output_path.endswith('.txt') else output_path + '.txt'
            with open(txt_output, 'w', encoding='utf-8') as f:
                f.writelines(result_lines)
            
            logger.info(f"Processed PDF ({len(pdf.pages)} pages) -> {txt_output}")
    except Exception as e:
        logger.error(f"Failed to process PDF '{input_path}': {e}")


def _process_html(input_path: str, output_path: str, config: PIRemoverConfig) -> None:
    """Process HTML file - extracts visible text, redacts PI."""
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        logger.error("beautifulsoup4 not installed. Run: pip install beautifulsoup4")
        return
    
    from pi_remover.remover import PIRemover
    
    try:
        with open(input_path, 'r', encoding='utf-8', errors='replace') as f:
            soup = BeautifulSoup(f.read(), 'html.parser')
        logger.info(f"Loaded HTML file: {input_path}")
    except Exception as e:
        logger.error(f"Failed to read HTML file '{input_path}': {e}")
        return
    
    remover = PIRemover(config)
    
    # Find all text nodes and redact
    for text_node in soup.find_all(string=True):
        if text_node.parent.name not in ['script', 'style', 'meta', 'link']:
            if text_node.strip():
                text_node.replace_with(remover.redact(str(text_node)))
    
    try:
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(str(soup))
        logger.info(f"Processed HTML -> {output_path}")
    except Exception as e:
        logger.error(f"Failed to write HTML output '{output_path}': {e}")


# Exports
__all__ = [
    # Validation
    'ValidationError',
    'validate_file',
    'validate_columns',
    'validate_config',
    # Processing functions
    'process_file',
    'process_csv',
    'process_dataframe',
]
